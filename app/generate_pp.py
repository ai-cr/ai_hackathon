
from io import BytesIO

from app.llm_logic.llm import PresentationOutput, PresentationPrompt
from app.llm_logic.unsplash_images import get_unsplash_image, get_image_object

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime as dt


# 16:9 dimensions
SLIDE_WIDTH  = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

IMG_WIDTH  = Inches(5.5)
MARGIN     = Inches(0.3)
FONT_SIZE  = Pt(15)


def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.strip().lstrip("#")
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def apply_background_color(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def get_scaled_height(image_bytes: BytesIO, target_width: Emu) -> Emu:
    """Calculate proportional height based on image aspect ratio and target width."""
    image_bytes.seek(0)
    with PILImage.open(image_bytes) as img:
        orig_width, orig_height = img.size
    aspect_ratio = orig_height / orig_width
    return Emu(int(target_width * aspect_ratio))


def add_slide_number(slide, slide_number: int, text_color: RGBColor):
    """Adds a small slide number in the bottom-right corner."""
    num_width  = Inches(0.6)
    num_height = Inches(0.3)
    left = SLIDE_WIDTH  - num_width  - MARGIN
    top  = SLIDE_HEIGHT - num_height - MARGIN

    box = slide.shapes.add_textbox(left, top, num_width, num_height)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_number)
    run = p.runs[0]
    run.font.size = Pt(10)
    run.font.color.rgb = text_color


def generate_ppt(prompt: PresentationOutput, presentation_prompt: PresentationPrompt = None, save_localy=False):
    filename = f"pp_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs = Presentation()

    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    theme = prompt.theme
    bg_color     = hex_to_rgb(theme.background_color_hex)
    accent_color = hex_to_rgb(theme.accent_color_hex)
    text_color   = hex_to_rgb(theme.text_color_hex)

    # --- Title Slide ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background_color(title_slide, bg_color)

    title_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = prompt.title
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = accent_color

    if presentation_prompt and presentation_prompt.optional_metadata:
        meta = presentation_prompt.optional_metadata
        meta_lines = []
        if meta.author_name:
            meta_lines.append(meta.author_name)
        if meta.company_organization:
            meta_lines.append(meta.company_organization)
        if meta.presentation_date:
            meta_lines.append(meta.presentation_date.strftime("%B %d, %Y"))

        if meta_lines:
            meta_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.0), Inches(1.5))
            tf_meta = meta_box.text_frame
            tf_meta.word_wrap = True
            for i, line in enumerate(meta_lines):
                p_meta = tf_meta.paragraphs[0] if i == 0 else tf_meta.add_paragraph()
                p_meta.text = line
                p_meta.alignment = PP_ALIGN.CENTER
                run_meta = p_meta.runs[0]
                run_meta.font.size = Pt(16)
                run_meta.font.color.rgb = text_color

    # --- Content Slides ---
    for slide_data in prompt.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background_color(slide, bg_color)

        has_image = bool(slide_data.image_query)

        if has_image:
            text_left  = IMG_WIDTH + MARGIN * 2
            text_width = SLIDE_WIDTH - text_left - MARGIN
        else:
            text_left  = MARGIN
            text_width = SLIDE_WIDTH - MARGIN * 2

        text_top = Inches(0.5)
        # Leave room at the bottom for slide number
        content_bottom = SLIDE_HEIGHT - Inches(0.5)

        # Slide title
        title_box = slide.shapes.add_textbox(text_left, text_top, text_width, Inches(0.9))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_data.slide_title
        run_title = p_title.runs[0]
        run_title.font.size = Pt(24)
        run_title.font.bold = True
        run_title.font.color.rgb = accent_color

        # Bullet points
        bullet_top    = text_top + Inches(1.0)
        bullet_height = content_bottom - bullet_top - Inches(0.1)
        bullet_box = slide.shapes.add_textbox(text_left, bullet_top, text_width, bullet_height)
        tf_bullets = bullet_box.text_frame
        tf_bullets.word_wrap = True
        for i, bullet in enumerate(slide_data.bullet_points):
            p = tf_bullets.paragraphs[0] if i == 0 else tf_bullets.add_paragraph()
            p.text = f"• {bullet}"
            run = p.runs[0]
            run.font.size = FONT_SIZE
            run.font.color.rgb = text_color
            # Empty line after each bullet
            p_space = tf_bullets.add_paragraph()
            p_space.text = ""

        # Speaker notes → only in the PowerPoint notes pane (visible in presenter view)
        if slide_data.speaker_notes:
            slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

        # Slide number bottom-right
        add_slide_number(slide, slide_data.slide_number, text_color)

        # Image (left side) – height scaled proportionally to width
        if has_image:
            image_url = get_unsplash_image(slide_data.image_query)
            if image_url:
                image_stream: BytesIO = get_image_object(image_url)
                img_height = get_scaled_height(image_stream, IMG_WIDTH)
                img_top = (content_bottom - img_height) // 2
                image_stream.seek(0)
                slide.shapes.add_picture(image_stream, MARGIN, img_top, IMG_WIDTH, img_height)

    if save_localy:
        prs.save(filename)
        return filename
    else:
        return prs